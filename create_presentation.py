from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 45, 85)
ACCENT_BLUE = RGBColor(65, 140, 220)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 245, 250)
TEXT_DARK = RGBColor(40, 40, 40)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = ACCENT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Add title text
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.space_before = Pt(10)
    title_para.space_after = Pt(10)
    title_frame.margin_left = Inches(0.5)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.line_spacing = 1.3
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, 
    "SMART ELECTION VOTING SYSTEM",
    "A Modern Solution for Secure & Fair Elections")

# Slide 2: Introduction
add_content_slide(prs, "Why We Need This System", [
    "✓ Traditional voting has issues: long queues, human errors, and fraud risks",
    "✓ Need for a modern, secure, and transparent voting process",
    "✓ Digital voting with face recognition ensures one person = one vote",
    "✓ Real-time results and complete voter transparency",
    "✓ Reduces administrative burden and increases trust in elections"
])

# Slide 3: What is Our System?
add_content_slide(prs, "Project Overview", [
    "🎯 A complete digital voting platform for elections",
    "🔐 Secure registration using facial recognition",
    "✅ Easy-to-use voting interface",
    "📊 Real-time election results and analytics",
    "👨‍💼 Admin dashboard for election management",
    "⚡ Works offline and online seamlessly"
])

# Slide 4: Key Features - Part 1
add_content_slide(prs, "Key Features (Part 1)", [
    "👤 Smart Registration: Facial recognition + personal details",
    "🔒 Biometric Security: Face-based identification prevents fraud",
    "💾 Secure Database: All voter data encrypted and protected",
    "🎯 Single Vote Protection: System ensures one vote per person",
    "📱 User-Friendly Interface: Simple to use for all age groups"
])

# Slide 5: Key Features - Part 2
add_content_slide(prs, "Key Features (Part 2)", [
    "📊 Live Results Dashboard: See election results in real-time",
    "👨‍💼 Admin Control: Full election management and monitoring",
    "🔔 Instant Notifications: System alerts for important events",
    "📈 Detailed Analytics: Understand voting patterns and trends",
    "🌐 Multi-Language Support: Accessible to diverse voters"
])

# Slide 6: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "🎨 Frontend: React with TypeScript (modern, fast UI)",
    "⚙️ Backend: Node.js & Express (reliable, scalable server)",
    "💾 Database: MongoDB (secure data storage)",
    "🔑 Security: JWT tokens & bcrypt encryption",
    "👁️ Face Recognition: face-api.js (works in browser, no server needed)",
    "📱 Responsive Design: Works on desktop, tablet, and mobile"
])

# Slide 7: How the System Works
add_content_slide(prs, "System Architecture", [
    "📲 User accesses the web application",
    "📝 User registers with face scan and personal details",
    "🎯 On election day, voter logs in and verifies identity",
    "🗳️ Selects preferred candidate from list",
    "✅ Confirms choice and vote is recorded",
    "📊 Admin sees real-time vote counts and results"
])

# Slide 8: Registration Process
add_content_slide(prs, "How Registration Works", [
    "Step 1️⃣: Fill in personal information (name, Aadhaar, etc.)",
    "Step 2️⃣: Allow camera access for face capture",
    "Step 3️⃣: System records facial features for later verification",
    "Step 4️⃣: Review all details and confirm registration",
    "Step 5️⃣: Get confirmation & unique voter token",
    "🔒 Data is securely stored and encrypted"
])

# Slide 9: Voting Process
add_content_slide(prs, "How Voting Works", [
    "Step 1️⃣: Log in with email and password",
    "Step 2️⃣: System verifies face matches registration record",
    "Step 3️⃣: View all candidates with their information",
    "Step 4️⃣: Select your preferred candidate",
    "Step 5️⃣: Review and confirm your vote",
    "✅ Vote is immediately locked and counted"
])

# Slide 10: Security Features
add_content_slide(prs, "Security & Safety", [
    "🔐 Password Encryption: Passwords are hashed with bcrypt",
    "👤 Face Recognition: Biometric verification prevents impersonation",
    "🔑 JWT Tokens: Secure session management",
    "📍 Role-Based Access: Admin and voter different permissions",
    "🚫 Fraud Prevention: One vote per person strictly enforced",
    "🔒 Data Protection: All sensitive data is encrypted"
])

# Slide 11: Admin Dashboard
add_content_slide(prs, "Admin Dashboard", [
    "📊 Election Statistics: Total votes, voters, candidates",
    "👥 Voter Management: View all registered voters",
    "🎯 Candidate Management: Add, edit, or remove candidates",
    "📈 Live Vote Tracking: See votes as they come in",
    "📋 Voter List: Download reports and analytics",
    "⚙️ System Settings: Configure election parameters"
])

# Slide 12: Results & Analytics
add_content_slide(prs, "Results & Winner Declaration", [
    "📊 Real-Time Dashboard: Live vote counting display",
    "🥇 Winner Announcement: Automatic winner calculation",
    "📈 Detailed Reports: Votes by candidate, pie charts, trends",
    "🔍 Transparency: Any voter can verify their vote was counted",
    "📁 Export Reports: Download results as PDF/CSV",
    "🎉 Celebration Mode: Special display for winner announcement"
])

# Slide 13: Future Enhancements
add_content_slide(prs, "What's Next? (Future Plans)", [
    "🌐 Blockchain Integration: For added transparency & security",
    "📱 Mobile App: Native apps for iOS and Android",
    "🌍 Multi-Regional: Support elections in multiple locations",
    "📧 Email/SMS: Send voting reminders to voters",
    "🎥 Live Streaming: Broadcast election results live",
    "🔊 Audio Support: Voice-based voting for accessibility"
])

# Slide 14: Conclusion
add_title_slide(prs,
    "LET'S BUILD THE FUTURE OF VOTING",
    "Together, we can make elections secure, fair, and transparent!")

# Save presentation
prs.save(r'd:\Mahesh\Code Library\election-voting-system\Election_Voting_System_Presentation.pptx')
print("✅ PowerPoint presentation created successfully!")
print("📁 File saved as: Election_Voting_System_Presentation.pptx")
